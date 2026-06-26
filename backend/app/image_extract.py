"""Extract images from PDF, DOCX, and PPTX files for vision ingest."""

from __future__ import annotations

import logging
import zipfile
from dataclasses import dataclass
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.config import settings

logger = logging.getLogger(__name__)

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff"}


@dataclass(frozen=True)
class ExtractedImage:
    label: str
    image_bytes: bytes
    mime_type: str


def _mime_for_ext(ext: str) -> str:
    ext = ext.lower().lstrip(".")
    if ext in {"jpg", "jpeg"}:
        return "image/jpeg"
    if ext == "png":
        return "image/png"
    if ext == "gif":
        return "image/gif"
    if ext in {"tif", "tiff"}:
        return "image/tiff"
    if ext == "bmp":
        return "image/bmp"
    if ext == "webp":
        return "image/webp"
    return "image/png"


def _large_enough(image_bytes: bytes) -> bool:
    return len(image_bytes) >= settings.VISION_MIN_IMAGE_BYTES


def extract_pdf_page_images(
    pdf_path: Path, page_num: int, page: fitz.Page, doc: fitz.Document
) -> list[ExtractedImage]:
    images: list[ExtractedImage] = []
    seen_hashes: set[str] = set()

    def add(label: str, image_bytes: bytes, mime_type: str) -> None:
        if not _large_enough(image_bytes):
            return
        digest = str(hash(image_bytes))
        if digest in seen_hashes:
            return
        seen_hashes.add(digest)
        images.append(ExtractedImage(label=label, image_bytes=image_bytes, mime_type=mime_type))

    text = (page.get_text() or "").strip()
    if len(text) < settings.VISION_MIN_TEXT_CHARS:
        pixmap = page.get_pixmap(dpi=settings.VISION_RENDER_DPI)
        add(
            f"{pdf_path.name} page {page_num} (full page)",
            pixmap.tobytes("png"),
            "image/png",
        )

    for img_index, image_info in enumerate(page.get_images(), start=1):
        if len(images) >= settings.VISION_MAX_IMAGES_PER_FILE:
            break
        try:
            xref = image_info[0]
            extracted = doc.extract_image(xref)
            add(
                f"{pdf_path.name} page {page_num} image {img_index}",
                extracted["image"],
                _mime_for_ext(str(extracted.get("ext", "png"))),
            )
        except Exception as exc:
            logger.debug(
                "Skipped PDF image %s p%s #%s: %s",
                pdf_path.name,
                page_num,
                img_index,
                exc,
            )

    return images


def extract_pdf_images(pdf_path: Path) -> list[ExtractedImage]:
    """Extract images from all pages (used for tests)."""
    images: list[ExtractedImage] = []
    try:
        doc = fitz.open(pdf_path)
    except Exception as exc:
        logger.warning("Could not open PDF %s for image extract: %s", pdf_path.name, exc)
        return images

    try:
        for page_num, page in enumerate(doc, start=1):
            if len(images) >= settings.VISION_MAX_IMAGES_PER_FILE:
                break
            for image in extract_pdf_page_images(pdf_path, page_num, page, doc):
                images.append(image)
                if len(images) >= settings.VISION_MAX_IMAGES_PER_FILE:
                    break
    finally:
        doc.close()

    return images


def extract_docx_images(docx_path: Path) -> list[ExtractedImage]:
    images: list[ExtractedImage] = []
    seen_hashes: set[str] = set()

    try:
        with zipfile.ZipFile(docx_path) as archive:
            media_names = sorted(
                name
                for name in archive.namelist()
                if name.startswith("word/media/")
                and Path(name).suffix.lower() in IMAGE_EXTENSIONS
            )
            for index, name in enumerate(media_names, start=1):
                if len(images) >= settings.VISION_MAX_IMAGES_PER_FILE:
                    break
                image_bytes = archive.read(name)
                if not _large_enough(image_bytes):
                    continue
                digest = str(hash(image_bytes))
                if digest in seen_hashes:
                    continue
                seen_hashes.add(digest)
                images.append(
                    ExtractedImage(
                        label=f"{docx_path.name} image {index} ({Path(name).name})",
                        image_bytes=image_bytes,
                        mime_type=_mime_for_ext(Path(name).suffix),
                    )
                )
    except Exception as exc:
        logger.warning("Could not extract images from DOCX %s: %s", docx_path.name, exc)

    return images


def extract_pptx_images(pptx_path: Path) -> list[ExtractedImage]:
    images: list[ExtractedImage] = []
    seen_hashes: set[str] = set()

    try:
        presentation = Presentation(str(pptx_path))
        for slide_num, slide in enumerate(presentation.slides, start=1):
            if len(images) >= settings.VISION_MAX_IMAGES_PER_FILE:
                break
            img_index = 0
            for shape in slide.shapes:
                if len(images) >= settings.VISION_MAX_IMAGES_PER_FILE:
                    break
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                try:
                    image_bytes = shape.image.blob
                    if not _large_enough(image_bytes):
                        continue
                    digest = str(hash(image_bytes))
                    if digest in seen_hashes:
                        continue
                    seen_hashes.add(digest)
                    img_index += 1
                    images.append(
                        ExtractedImage(
                            label=f"{pptx_path.name} slide {slide_num} image {img_index}",
                            image_bytes=image_bytes,
                            mime_type=_mime_for_ext(shape.image.ext),
                        )
                    )
                except Exception as exc:
                    logger.debug(
                        "Skipped PPTX image %s slide %s: %s",
                        pptx_path.name,
                        slide_num,
                        exc,
                    )
    except Exception as exc:
        logger.warning("Could not extract images from PPTX %s: %s", pptx_path.name, exc)

    return images
